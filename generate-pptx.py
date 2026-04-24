#!/usr/bin/env python3
"""Generate a .pptx slide deck: The Semantic Layer as the Enterprise Knowledge Graph"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Colors
BG = RGBColor(0x1A, 0x1A, 0x2E)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
RED = RGBColor(0xFF, 0x6B, 0x6B)
WHITE = RGBColor(0xE0, 0xE0, 0xE0)
DARK_BLUE = RGBColor(0x0F, 0x34, 0x60)
GRAY = RGBColor(0x90, 0x90, 0x90)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

def fmt(run, size=18, color=WHITE, bold=False, italic=False, font_name='Helvetica Neue'):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)
    tb = add_textbox(slide, 1, 1.5, 11, 2)
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    fmt(r, 44, CYAN, bold=True)

    p2 = tb.text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)
    r2 = p2.add_run()
    r2.text = subtitle
    fmt(r2, 24, WHITE)
    return slide

def add_content_slide(title, bullets, subnote=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    tb = add_textbox(slide, 0.8, 0.4, 11, 1)
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    fmt(r, 36, CYAN, bold=True)

    tb2 = add_textbox(slide, 1.0, 1.5, 11, 5)
    tb2.text_frame.word_wrap = True
    first = True
    for bullet in bullets:
        if first:
            p = tb2.text_frame.paragraphs[0]
            first = False
        else:
            p = tb2.text_frame.add_paragraph()
        p.space_before = Pt(10)
        p.space_after = Pt(4)

        if isinstance(bullet, tuple):
            # (text, color, bold, size)
            text, color, bold, size = bullet
            r = p.add_run()
            r.text = text
            fmt(r, size, color, bold)
        else:
            r = p.add_run()
            r.text = bullet
            fmt(r, 20, WHITE)

    if subnote:
        p = tb2.text_frame.add_paragraph()
        p.space_before = Pt(20)
        r = p.add_run()
        r.text = subnote
        fmt(r, 16, GRAY, italic=True)

    return slide

def add_two_column_slide(title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    tb = add_textbox(slide, 0.8, 0.4, 11, 1)
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    fmt(r, 36, CYAN, bold=True)

    # Left
    ltb = add_textbox(slide, 0.8, 1.6, 5.5, 5)
    ltb.text_frame.word_wrap = True
    p = ltb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = left_title
    fmt(r, 24, CYAN, bold=True)
    for item in left_items:
        p = ltb.text_frame.add_paragraph()
        p.space_before = Pt(8)
        r = p.add_run()
        r.text = item
        fmt(r, 18, WHITE)

    # Right
    rtb = add_textbox(slide, 7.0, 1.6, 5.5, 5)
    rtb.text_frame.word_wrap = True
    p = rtb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = right_title
    fmt(r, 24, CYAN, bold=True)
    for item in right_items:
        p = rtb.text_frame.add_paragraph()
        p.space_before = Pt(8)
        r = p.add_run()
        r.text = item
        fmt(r, 18, WHITE)

    return slide

def add_table_slide(title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    tb = add_textbox(slide, 0.8, 0.4, 11, 1)
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    fmt(r, 36, CYAN, bold=True)

    cols = len(headers)
    row_count = len(rows) + 1
    tbl_shape = slide.shapes.add_table(row_count, cols, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5 * row_count + 0.3))
    tbl = tbl_shape.table

    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                fmt(run, 16, CYAN, bold=True)
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x16, 0x21, 0x3E)

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    fmt(run, 14, WHITE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE

    return slide

# ── Slides ──

# 1 - Title
s = add_title_slide(
    "You Don't Need a Knowledge Graph",
    "You Already Have One"
)
tb = add_textbox(s, 1, 4.2, 11, 1)
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "The Semantic Layer as the Enterprise Knowledge Graph"
fmt(r, 20, RED, bold=True)

# 2 - The Problem
add_content_slide("The Problem", [
    "LLMs are great at understanding natural language.",
    "",
    "They are terrible at:",
    "    \u2022  Being databases",
    "    \u2022  Doing math",
    "    \u2022  Remembering facts",
    "    \u2022  Statistical reasoning",
], subnote="The goal: make them use actual facts from data, not statistical correlation.")

# 3 - Conventional Wisdom
add_content_slide("The Conventional Wisdom", [
    ('"You need a Knowledge Graph for AI"', WHITE, True, 24),
    "",
    "Build a Neo4j database. ETL your warehouse into it.",
    "Nodes, edges, Cypher queries. A whole new infrastructure layer.",
    "",
    ("But where do the relationships already live?", CYAN, True, 22),
    "",
    "In your relational database.",
    "Defined by foreign keys, joins, and business logic.",
])

# 4 - The Question
add_title_slide(
    "Two Paths to the Same Answer",
    '"Who issued John Smith\'s insurance policy?"'
)

# 5 - Path 1
add_content_slide("Path 1: Knowledge Graph + Neo4j", [
    "1.  ETL relational data into Neo4j",
    "2.  Create nodes: Agent, Policy, Customer",
    "3.  Create edges: ISSUED, COVERS",
    "4.  LLM generates Cypher:",
    "",
    '    MATCH (a:Agent)-[:ISSUED]->(p:Policy)',
    '          -[:COVERS]->(c {name:"John Smith"})',
    '    RETURN a.name',
    "",
    ("Cost: Duplicate data store, ETL pipeline, new infrastructure", RED, True, 18),
])

# 6 - Path 2
add_content_slide("Path 2: Semantic Layer + SQL", [
    "1.  Relationships already exist in the warehouse",
    "2.  SML defines: Policy \u2192 Agent via agent_id, Policy \u2192 Customer via customer_id",
    "3.  LLM reads SML, generates SQL:",
    "",
    "    SELECT a.name",
    "    FROM policy p",
    "    JOIN agent a ON p.agent_id = a.id",
    "    JOIN customer c ON p.customer_id = c.id",
    "    WHERE c.name = 'John Smith'",
    "",
    ("Cost: None. Data stays where it is.", CYAN, True, 18),
])

# 7 - Comparison Table
add_table_slide("Side by Side", [
    "", "Knowledge Graph", "Semantic Layer"
], [
    ["Relationships", "Explicit edges in graph DB", "Join definitions in SML"],
    ["Data location", "Duplicated into Neo4j", "Stays in warehouse"],
    ["Query language", "Cypher", "SQL (generated)"],
    ["Extra infra", "Neo4j + ETL pipeline", "None"],
    ["Data freshness", "Depends on ETL lag", "Real-time"],
    ["Scales with", "Separate concern", "Inherits warehouse scale"],
])

# 8 - Complex Traversals
add_content_slide('"But What About Complex Traversals?"', [
    "The classic Neo4j argument:",
    "",
    ('"Find everyone within 3 degrees of John Smith"', WHITE, False, 22),
    "",
    "That requires a complex Cypher query...",
    "",
    ("...or three simple SQL queries.", CYAN, True, 24),
])

# 9 - Decomposition
add_content_slide("LLMs Don't Need One Perfect Query", [
    "An LLM with agentic reasoning decomposes:",
    "",
    ("Step 1:", CYAN, True, 20),
    '    "Who issued John Smith\'s policy?" \u2192 Agent ID',
    "",
    ("Step 2:", CYAN, True, 20),
    '    "What other policies did that agent issue?" \u2192 Policy list',
    "",
    ("Step 3:", CYAN, True, 20),
    '    "Who are the customers on those policies?" \u2192 Names',
    "",
    "Three trivial, auditable, governed SQL queries.",
])

# 10 - Why Decomposition Wins
add_two_column_slide(
    "Why Decomposition Wins",
    "Complex Cypher (one shot)",
    [
        "MATCH (a)-[:ISSUED]->(p)",
        "  -[:COVERS]->(c)",
        "  -[:RELATED_TO*1..3]-(x)",
        "WHERE c.name = 'John Smith'",
        "RETURN DISTINCT x.name",
        "",
        "\u2022 Single opaque operation",
        "\u2022 Hard to debug if wrong",
    ],
    "Multi-step SQL (agentic)",
    [
        "\u2022 Each step is auditable",
        "\u2022 Each step is simple SQL",
        "\u2022 Less likely to be wrong",
        "\u2022 Semantic layer governs every step",
        "\u2022 Consistent metric definitions",
        "\u2022 Access control enforced",
    ]
)

# 11 - Error Modes
add_content_slide("What Actually Causes LLM Errors?", [
    "The failure modes are:",
    "",
    "    \u2022  Hallucinating a column name that doesn't exist",
    "    \u2022  Joining on the wrong key",
    "    \u2022  Misunderstanding what a metric means",
    "    \u2022  Missing a filter (soft deletes, date ranges)",
    "",
    ("The semantic layer addresses ALL of these.", CYAN, True, 22),
    "",
    ("A graph database addresses NONE of them.", RED, True, 22),
    "",
    "You can hallucinate Cypher just as easily as SQL.",
])

# 12 - The Principle
s = add_title_slide(
    "The Principle",
    ""
)
tb = add_textbox(s, 1, 3.5, 11, 2)
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "LLMs should be the natural language interface,\nnot the computation engine."
fmt(r, 28, WHITE)

# 13 - Right Tool
add_table_slide("The Right Tool for Each Job", [
    "Task", "LLM Role", "Delegate To"
], [
    ['"15% of $4,230?"', "Parse question", "Calculator"],
    ['"Who issued the policy?"', "Understand intent", "Semantic layer + SQL"],
    ['"Is this fraudulent?"', "Frame question", "Bayesian model"],
    ['"Last quarter\'s revenue?"', "Present answer", "Governed metrics"],
])

# 14 - Over Time
add_content_slide("The Neo4j Case Gets Weaker Over Time", [
    "As LLMs improve at multi-step reasoning:",
    "",
    '    \u2022  "Complex query" advantage of graph DBs erodes',
    "       LLMs decompose complex questions into simple ones",
    "",
    '    \u2022  "Schema understanding" advantage of semantic layers grows',
    "       LLMs read SML and reason about relationships",
    "",
    "    \u2022  Data duplication cost stays constant or worsens",
    "       Data volumes only go up",
])

# 15 - When You Need Graph
add_content_slide("When You Actually Need a Graph Database", [
    "    \u2022  Data is natively graph-shaped (social networks, fraud rings)",
    "    \u2022  You need graph algorithms (PageRank, shortest path)",
    "    \u2022  Integrating unstructured sources (entity extraction)",
    "",
    "",
    "For enterprise analytics against warehouse data?",
    "",
    ("The semantic layer IS the knowledge graph.", CYAN, True, 28),
])

# 16 - Summary
add_content_slide("Summary", [
    ("1.", CYAN, True, 22),
    "   The semantic layer already encodes the relationship",
    "   metadata a knowledge graph would duplicate",
    "",
    ("2.", CYAN, True, 22),
    "   LLMs + simple tools beat LLMs + complex tools",
    "",
    ("3.", CYAN, True, 22),
    "   Use LLMs for language, delegate everything else",
    "   to authoritative tools",
    "",
    ("4.", CYAN, True, 22),
    "   System quality depends on tool quality,",
    "   not on making the LLM smarter",
])

# 17 - Thank You
s = add_title_slide(
    "Thank You",
    "The Semantic Layer as the Enterprise Knowledge Graph"
)
tb = add_textbox(s, 1, 4.5, 11, 1)
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "LLMs for language. Real tools for real answers."
fmt(r, 20, GRAY, italic=True)


out = "/Users/chuck/atscale/slides/deck.pptx"
prs.save(out)
print(f"Saved: {out}")
